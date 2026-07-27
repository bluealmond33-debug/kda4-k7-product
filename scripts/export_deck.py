#!/usr/bin/env python3
"""KARI-NA 발표 덱 내보내기 — docs/deck.html → PDF · PPTX.

덱의 `?shot=N` 모드로 각 장을 크롬(레일·목차·노트) 없이 1280×720 원본 크기로 띄운 뒤,
헤드리스 크롬으로 2배 해상도(2560×1440) 캡처한다. 그 이미지를 그대로
PDF 페이지와 16:9 PPTX 슬라이드로 싣는다 — 화면과 픽셀 단위로 같은 결과가 나온다.

    python3 scripts/export_deck.py

산출: docs/export/KARI-NA-발표.pdf · KARI-NA-발표.pptx · slides/*.png
"""

import re
import shutil
import subprocess
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
DECK = ROOT / "docs" / "deck.html"
OUT = ROOT / "docs" / "export"
SHOTS = OUT / "slides"

CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
SCALE = 2                      # 2560×1440 캡처
W, H = 1280, 720               # 덱의 무대 크기 (16:9)


def slide_count() -> int:
    return len(re.findall(r'<section class="slide', DECK.read_text(encoding="utf-8")))


def capture(n: int) -> Path:
    """n번 장을 PNG로 캡처한다 (1-indexed)."""
    png = SHOTS / f"{n:02d}.png"
    subprocess.run(
        [
            CHROME,
            "--headless=new",
            "--disable-gpu",
            "--hide-scrollbars",
            "--virtual-time-budget=3000",
            f"--force-device-scale-factor={SCALE}",
            f"--window-size={W},{H}",
            f"--screenshot={png}",
            f"file://{DECK}?shot={n}",
        ],
        check=True,
        capture_output=True,
    )
    if not png.exists():
        raise RuntimeError(f"{n}번 장 캡처 실패")
    return png


def build_pdf(pngs: list[Path], dest: Path) -> None:
    pages = [Image.open(p).convert("RGB") for p in pngs]
    pages[0].save(dest, save_all=True, append_images=pages[1:], resolution=144.0)


def build_pptx(pngs: list[Path], dest: Path) -> None:
    prs = Presentation()
    # 16:9 — 13.333in × 7.5in
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(png), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
    prs.save(dest)


def main() -> int:
    if not Path(CHROME).exists():
        print(f"크롬을 찾을 수 없습니다: {CHROME}", file=sys.stderr)
        return 1

    total = slide_count()
    if SHOTS.exists():
        shutil.rmtree(SHOTS)
    SHOTS.mkdir(parents=True)

    pngs = []
    for n in range(1, total + 1):
        pngs.append(capture(n))
        print(f"  캡처 {n:2d}/{total}", end="\r", flush=True)
    print(f"  캡처 {total}/{total} 완료      ")

    pdf = OUT / "KARI-NA-발표.pdf"
    pptx = OUT / "KARI-NA-발표.pptx"
    build_pdf(pngs, pdf)
    build_pptx(pngs, pptx)

    for f in (pdf, pptx):
        print(f"  {f.relative_to(ROOT)}  ({f.stat().st_size / 1024:.0f} KB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
